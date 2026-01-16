#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create Comprehensive PPT Template

Creates a realistic PowerPoint template with:
- Slide 1: summary dashboard (business-style conclusion cards)
- Slide 2: overview (title, 4 charts, key findings)
- Slide 3: detailed metrics (tables + notes)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
import os


# ==================== STYLE CONSTANTS ====================
# Unified style configuration for consistent formatting across slides
# McKinsey-inspired professional color scheme

# Colors - McKinsey Style
COLOR_PRIMARY = RGBColor(36, 71, 127)      # McKinsey deep blue (#24477f)
COLOR_TITLE = RGBColor(36, 71, 127)        # Deep blue for titles
COLOR_SUBTITLE = RGBColor(100, 100, 100)   # Gray for subtitles
COLOR_BODY = RGBColor(50, 50, 50)          # Dark gray for body text
COLOR_CHART_PRIMARY = RGBColor(36, 71, 127)  # Primary chart color (McKinsey blue)
COLOR_TABLE_HEADER_BG = RGBColor(36, 71, 127)  # Table header background
COLOR_TABLE_HEADER_TEXT = RGBColor(255, 255, 255)  # Table header text (white)

# Font sizes
FONT_SIZE_TITLE = 24            # Main slide title
FONT_SIZE_SUBTITLE = 12         # Subtitle text
FONT_SIZE_BODY = 11             # Body text, conclusions, notes
FONT_SIZE_CHART_TITLE = 14      # Chart titles
FONT_SIZE_CHART_AXIS = 9        # Chart axis labels (tick labels)
FONT_SIZE_CHART_LEGEND = 9      # Chart legend text
FONT_SIZE_CHART_DATA_LABEL = 8  # Chart data labels (values on points)
FONT_SIZE_TABLE_HEADER = 11     # Table headers
FONT_SIZE_TABLE_BODY = 9        # Table body text

# Layout positions
TITLE_LEFT = Inches(0.5)
TITLE_TOP = Inches(0.3)
TITLE_WIDTH = Inches(9)
TITLE_HEIGHT = Inches(0.5)

SUBTITLE_LEFT = Inches(0.5)
SUBTITLE_TOP = Inches(0.85)
SUBTITLE_WIDTH = Inches(9)
SUBTITLE_HEIGHT = Inches(0.35)

CONTENT_TOP = Inches(1.4)  # Unified starting position for charts/tables
# ========================================================


def create_text_shape(slide, left, top, width, height, text,
                      font_size=14, bold=False, color=None):
    """Helper to create a text box."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color

    return text_box


def create_line_chart(slide, left, top, width, height, title):
    """Create a line chart with sample data."""
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('Series 1', [0.75, 0.78, 0.81, 0.84])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        left,
        top,
        width,
        height,
        chart_data,
    )

    chart = chart_shape.chart

    # Chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(FONT_SIZE_CHART_TITLE)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # Apply McKinsey-style single color to series
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = COLOR_CHART_PRIMARY
    # Set line color for line chart
    series.format.line.color.rgb = COLOR_CHART_PRIMARY
    series.format.line.width = Pt(2)

    # Legend styling
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(FONT_SIZE_CHART_LEGEND)  # Keep legend subtle

    # Axis formatting - for 0.82 -> 82%
    value_axis_labels = chart.value_axis.tick_labels
    value_axis_labels.number_format = '0%'
    value_axis_labels.number_format_is_linked = False
    value_axis_labels.font.size = Pt(FONT_SIZE_CHART_AXIS)  # Smaller for less distraction
    chart.category_axis.tick_labels.font.size = Pt(FONT_SIZE_CHART_AXIS)

    # Data labels - show values directly on data points (0.82 -> 82%)
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '0%'
    data_labels.number_format_is_linked = False
    data_labels.font.size = Pt(FONT_SIZE_CHART_DATA_LABEL)  # Small but readable

    # Remove gridlines for cleaner appearance
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False

    return chart_shape


def create_bar_chart(slide, left, top, width, height, title):
    """Create a bar chart with sample data."""
    chart_data = CategoryChartData()
    chart_data.categories = ['18-25', '26-35', '36-45', '46-55', '56+']
    chart_data.add_series('Distribution', [20, 30, 25, 15, 10])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left,
        top,
        width,
        height,
        chart_data,
    )

    chart = chart_shape.chart

    # Chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(FONT_SIZE_CHART_TITLE)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # Apply McKinsey-style single color to bar series
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = COLOR_CHART_PRIMARY

    # No legend for bar chart
    chart.has_legend = False

    # Axis formatting - use custom format to show % without multiplying by 100
    value_axis_labels = chart.value_axis.tick_labels
    value_axis_labels.number_format = '0"%"'
    value_axis_labels.number_format_is_linked = False
    value_axis_labels.font.size = Pt(FONT_SIZE_CHART_AXIS)
    chart.category_axis.tick_labels.font.size = Pt(FONT_SIZE_CHART_AXIS)

    # Data labels - use custom format to show % without multiplying by 100
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '0.0"%"'
    data_labels.number_format_is_linked = False
    data_labels.font.size = Pt(FONT_SIZE_CHART_DATA_LABEL)

    # Remove gridlines for cleaner appearance
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False

    return chart_shape


def create_report_slide(prs, department_name="Department Name"):
    """Create a comprehensive overview slide (page 2)."""

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title section
    title_box = create_text_shape(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        f"{department_name} - Quarterly Survey Report",
        font_size=FONT_SIZE_TITLE,
        bold=True,
        color=COLOR_TITLE,
    )
    title_box.name = "TextBox_Title"

    # Subtitle with key metrics
    subtitle_box = create_text_shape(
        slide,
        SUBTITLE_LEFT,
        SUBTITLE_TOP,
        SUBTITLE_WIDTH,
        SUBTITLE_HEIGHT,
        "Q4 2024 | Department Size: XXX | Response Rate: XX%",
        font_size=FONT_SIZE_SUBTITLE,
        color=COLOR_SUBTITLE,
    )
    subtitle_box.name = "TextBox_Subtitle"

    # Four charts in 2x2 grid
    chart_width = Inches(4.2)
    chart_height = Inches(2.6)
    margin_left = Inches(0.5)
    margin_top = CONTENT_TOP  # Use unified constant
    h_spacing = Inches(0.3)
    v_spacing = Inches(0.2)

    # Top-left: Satisfaction Trend
    chart_satisfaction = create_line_chart(
        slide,
        margin_left,
        margin_top,
        chart_width,
        chart_height,
        "Satisfaction",
    )
    chart_satisfaction.name = "Chart_Satisfaction"

    # Top-right: Engagement Trend
    chart_engagement = create_line_chart(
        slide,
        margin_left + chart_width + h_spacing,
        margin_top,
        chart_width,
        chart_height,
        "Engagement",
    )
    chart_engagement.name = "Chart_Engagement"

    # Bottom-left: Response Rate Trend
    chart_response = create_line_chart(
        slide,
        margin_left,
        margin_top + chart_height + v_spacing,
        chart_width,
        chart_height,
        "Response Rate",
    )
    chart_response.name = "Chart_Response"

    # Bottom-right: Age Distribution (Bar chart)
    chart_age = create_bar_chart(
        slide,
        margin_left + chart_width + h_spacing,
        margin_top + chart_height + v_spacing,
        chart_width,
        chart_height,
        "Age Distribution",
    )
    chart_age.name = "Chart_Age"

    # Conclusion section at bottom
    conclusion_top = margin_top + 2 * chart_height + 2 * v_spacing + Inches(0.15)
    conclusion_box = create_text_shape(
        slide,
        Inches(0.5),
        conclusion_top,
        Inches(9),
        Inches(0.6),
        "Key Findings: Overall satisfaction increased by X% QoQ. "
        "Engagement remains stable with strong performance in 26-35 age group.",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    conclusion_box.name = "TextBox_Conclusion"

    print(f"  ✓ Created overview slide: {department_name}")


def create_detail_slide(prs, department_name="[Department Name]"):
    """Create a detailed metrics slide (page 3) with tables."""

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = create_text_shape(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        f"{department_name} - Detailed Metrics",
        font_size=FONT_SIZE_TITLE,  # Now consistent with slide 1 (24pt)
        bold=True,
        color=COLOR_TITLE,
    )
    title_box.name = "TextBox_DetailTitle"

    # Subtitle
    subtitle_box = create_text_shape(
        slide,
        SUBTITLE_LEFT,
        SUBTITLE_TOP,
        SUBTITLE_WIDTH,
        SUBTITLE_HEIGHT,
        "QoQ changes and demographic breakdown",
        font_size=FONT_SIZE_SUBTITLE,
        color=COLOR_SUBTITLE,
    )
    subtitle_box.name = "TextBox_DetailSubtitle"

    # Key metrics table (left side)
    rows_metrics = 5  # 1 header + 4 metrics
    cols_metrics = 4  # Metric, Q3, Q4, QoQ Δ
    metrics_left = Inches(0.5)
    metrics_top = CONTENT_TOP  # Now consistent with slide 1 (1.4 inches)
    metrics_width = Inches(5.2)
    metrics_height = Inches(3.0)

    metrics_shape = slide.shapes.add_table(
        rows_metrics,
        cols_metrics,
        metrics_left,
        metrics_top,
        metrics_width,
        metrics_height,
    )
    metrics_shape.name = "Table_Metrics"
    metrics_table = metrics_shape.table

    # Header row - McKinsey style (deep blue background + white text)
    headers = ["Metric", "Q3", "Q4", "QoQ Δ"]
    for col_idx, header in enumerate(headers):
        cell = metrics_table.cell(0, col_idx)
        # Apply McKinsey blue background
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER_BG
        # Set text
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(FONT_SIZE_TABLE_HEADER)
            paragraph.font.color.rgb = COLOR_TABLE_HEADER_TEXT  # White text

    # Default metric names (data will overwrite in automation)
    default_metrics = ["Satisfaction", "Engagement", "Response Rate", "Retention"]
    for row_idx, metric_name in enumerate(default_metrics, start=1):
        # Set all cells in this row with proper formatting
        for col_idx in range(cols_metrics):
            cell = metrics_table.cell(row_idx, col_idx)
            text_frame = cell.text_frame
            text_frame.clear()  # Clear any default content
            p = text_frame.paragraphs[0]
            run = p.add_run()

            # Set initial text
            if col_idx == 0:
                run.text = metric_name
            else:
                run.text = "--"  # Placeholder that will be replaced

            # Set font properties on the run to ensure format preservation
            run.font.size = Pt(FONT_SIZE_TABLE_BODY)

    # Age distribution table (right side)
    rows_age = 6  # 1 header + 5 age groups
    cols_age = 2  # Age group, share
    age_left = Inches(5.9)
    age_top = CONTENT_TOP  # Now consistent with metrics table (1.4 inches)
    age_width = Inches(3.6)
    age_height = Inches(3.0)

    age_shape = slide.shapes.add_table(
        rows_age,
        cols_age,
        age_left,
        age_top,
        age_width,
        age_height,
    )
    age_shape.name = "Table_Age"
    age_table = age_shape.table

    # Header row - McKinsey style (deep blue background + white text)
    age_headers = ["Age Group", "Share of respondents"]
    for col_idx, header in enumerate(age_headers):
        cell = age_table.cell(0, col_idx)
        # Apply McKinsey blue background
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER_BG
        # Set text
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(FONT_SIZE_TABLE_HEADER)
            paragraph.font.color.rgb = COLOR_TABLE_HEADER_TEXT  # White text

    # Default age groups (data will overwrite in automation)
    default_age_groups = ["18-25", "26-35", "36-45", "46-55", "56+"]
    for row_idx, age_group in enumerate(default_age_groups, start=1):
        # Set all cells in this row with proper formatting
        for col_idx in range(cols_age):
            cell = age_table.cell(row_idx, col_idx)
            text_frame = cell.text_frame
            text_frame.clear()  # Clear any default content
            p = text_frame.paragraphs[0]
            run = p.add_run()

            # Set initial text
            if col_idx == 0:
                run.text = age_group
            else:
                run.text = "--"  # Placeholder that will be replaced

            # Set font properties on the run to ensure format preservation
            run.font.size = Pt(FONT_SIZE_TABLE_BODY)

    # Key Insights area at bottom (consistent with page 1 style)
    notes_top = metrics_top + metrics_height + Inches(0.3)
    notes_box = create_text_shape(
        slide,
        Inches(0.5),
        notes_top,
        Inches(9),
        Inches(0.6),
        "Key Insights: Satisfaction improved by X pp from Q3 to Q4. "
        "Largest respondent group: 26-35 age range.",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    notes_box.name = "TextBox_Notes"

    print(f"  ✓ Created detail slide: {department_name}")


def create_dashboard_slide(prs, department_name="[Department Name]"):
    """Create a business-style summary dashboard slide (page 1).

    This slide uses a multi-card layout to present high-level conclusions
    across key survey metrics: satisfaction, engagement, response rate,

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = create_text_shape(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        f"{department_name} - Summary Dashboard",
        font_size=FONT_SIZE_TITLE,
        bold=True,
        color=COLOR_TITLE,
    )
    title_box.name = "TextBox_DashboardTitle"

    # Subtitle
    subtitle_box = create_text_shape(
        slide,
        SUBTITLE_LEFT,
        SUBTITLE_TOP,
        SUBTITLE_WIDTH,
        SUBTITLE_HEIGHT,
        "Key conclusions across satisfaction, engagement, response, retention and age structure",
        font_size=FONT_SIZE_SUBTITLE,
        color=COLOR_SUBTITLE,
    )
    subtitle_box.name = "TextBox_DashboardSubtitle"

    # Card layout: 2 rows × 3 columns
    card_width = Inches(3.0)
    card_height = Inches(1.4)
    margin_left = Inches(0.5)
    margin_top = CONTENT_TOP
    h_spacing = Inches(0.25)
    v_spacing = Inches(0.35)

    # Row 1
    card1 = create_text_shape(
        slide,
        margin_left,
        margin_top,
        card_width,
        card_height,
        "Conclusion 1:\nOverall Satisfaction",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    card1.name = "Card_OverallSatisfaction"

    card2 = create_text_shape(
        slide,
        margin_left + card_width + h_spacing,
        margin_top,
        card_width,
        card_height,
        "Conclusion 2:\nEmployee Engagement",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    card2.name = "Card_Engagement"

    card3 = create_text_shape(
        slide,
        margin_left + 2 * (card_width + h_spacing),
        margin_top,
        card_width,
        card_height,
        "Conclusion 3:\nResponse & Retention",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    card3.name = "Card_ResponseRetention"

    # Row 2
    row2_top = margin_top + card_height + v_spacing

    card4 = create_text_shape(
        slide,
        margin_left,
        row2_top,
        card_width,
        card_height,
        "Conclusion 4:\nAge Structure",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    card4.name = "Card_AgeStructure"

    card5 = create_text_shape(
        slide,
        margin_left + card_width + h_spacing,
        row2_top,
        card_width,
        card_height,
        "Conclusion 5:\nSatisfaction vs Engagement",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    card5.name = "Card_SatVsEngagement"

    card6 = create_text_shape(
        slide,
        margin_left + 2 * (card_width + h_spacing),
        row2_top,
        card_width,
        card_height,
        "Conclusion 6:\nOverall Summary",
        font_size=FONT_SIZE_BODY,
        color=COLOR_BODY,
    )
    card6.name = "Card_OverallSummary"

    print(f"  ✓ Created dashboard slide: {department_name}")


def create_template():
    """Create comprehensive PPT template with dashboard, overview and detail slides."""

    print("\n" + "=" * 80)
    print("📊 Creating Comprehensive PPT Template")
    print("=" * 80)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create dashboard, overview and detail slides (will be used for all departments)
    create_dashboard_slide(prs, "[Department Name]")
    create_report_slide(prs, "[Department Name]")
    create_detail_slide(prs, "[Department Name]")

    # Save template
    output_path = 'data/raw/template.pptx'
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print("\n" + "=" * 80)
    print("✅ Template Created Successfully")
    print("=" * 80)
    print(f"\nOutput: {output_path}")
    print(f"Slides: {len(prs.slides)}")
    print("\nTemplate features:")
    print("  - Slide 1: Summary dashboard (business-style conclusion cards)")
    print("    • 6 cards summarizing satisfaction, engagement, response, retention, age structure and overall summary")
    print("  - Slide 2: Overview")
    print("    • Title with department name (dynamic)")
    print("    • Subtitle with key metrics (dynamic)")
    print("    • 4 charts in 2×2 grid:")
    print("      - Satisfaction trend (line chart)")
    print("      - Engagement trend (line chart)")
    print("      - Response rate trend (line chart)")
    print("      - Age distribution (bar chart)")
    print("    • Key findings section (dynamic)")
    print("  - Slide 3: Detailed metrics")
    print("    • Title and subtitle")
    print("    • Key metrics QoQ table (Satisfaction, Engagement, Response Rate, Retention)")
    print("    • Age distribution table")
    print("    • Notes / actions area")
    print("\n💡 Next Steps:")
    print("  1. Run: python batch_generate.py")
    print("  2. Check: data/output/reports/")
    print("=" * 80 + "\n")


if __name__ == '__main__':
    create_template()
