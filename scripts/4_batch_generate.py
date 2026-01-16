#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Batch PPT Generation Script

Demonstrates the power of automation by generating multiple reports
from a single template in seconds.
"""

import os
import sys

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData

# Add project root to path
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from src.mapping import load_mapping

# Load data→PPT mapping from markdown spec (documentation-as-code pattern)
_MAPPING_CONFIG = load_mapping()
SHAPE_NAMES = _MAPPING_CONFIG["shape_names"]
CHART_MAPPINGS = [
    # (logical_name, title_keyword, data_key, number_format)
    (m.logical_name, m.title_keyword or "", m.data_key, m.number_format)
    for m in _MAPPING_CONFIG["charts"]
    if m.slide_index == 0
]


def update_text_preserve_format(shape, new_text):
    """Update text while preserving formatting"""
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        text_frame.paragraphs[0].runs[0].text = new_text
    else:
        text_frame.text = new_text


def update_table_cell_preserve_format(cell, new_text):
    """
    Update table cell text while preserving formatting

    CRITICAL: Using cell.text = ... will DESTROY all formatting.
    Instead, we update the first run's text to preserve font size, color, etc.

    Args:
        cell: Table cell object
        new_text: New text content
    """
    text_frame = cell.text_frame
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        # Preserve formatting by updating the first run's text only
        text_frame.paragraphs[0].runs[0].text = new_text
    else:
        # Fallback: this may lose formatting
        cell.text = new_text


def update_chart_data(
    chart, categories, values, number_format=None, auto_scale=True, scale_factor=1.2
):
    """
    Update chart data while preserving template styling

    Args:
        chart: PPT chart object
        categories: X-axis categories
        values: Data values
        number_format: Number format (e.g., '0%')
        auto_scale: Enable Y-axis auto-scaling
        scale_factor: Scale factor for max value (default 1.2 = 120% of max)
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    # Add series (use existing series name from template if available)
    if chart.series:
        series_name = chart.series[0].name
    else:
        series_name = "Series 1"

    chart_data.add_series(series_name, values)
    chart.replace_data(chart_data)

    # CRITICAL: Re-apply formatting after replace_data() which may reset some properties
    try:
        from pptx.util import Pt

        # Re-apply axis label font sizes and number format
        if hasattr(chart, 'category_axis') and chart.category_axis:
            chart.category_axis.tick_labels.font.size = Pt(9)
        if hasattr(chart, 'value_axis') and chart.value_axis:
            value_labels = chart.value_axis.tick_labels
            value_labels.font.size = Pt(9)
            if number_format:
                value_labels.number_format = number_format
                value_labels.number_format_is_linked = False

        # Re-apply legend font size
        if hasattr(chart, 'legend') and chart.legend:
            chart.legend.font.size = Pt(9)

        # Re-apply data labels
        if len(chart.plots) > 0:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(8)
            if number_format:
                data_labels.number_format = number_format
                data_labels.number_format_is_linked = False
    except Exception as e:
        # If formatting fails, at least data is correct
        pass

    # Auto-scale Y-axis based on data range
    if auto_scale and values:
        # Filter out None values
        valid_values = [v for v in values if v is not None]

        if valid_values:
            min_val = min(valid_values)
            max_val = max(valid_values)

            # Set axis range with padding
            y_min = max(0, min_val * 0.9)  # 90% of min, but not below 0
            y_max = max_val * scale_factor  # 120% of max by default

            try:
                from pptx.oxml import parse_xml

                chart_xml = chart._element
                namespaces = {
                    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart"
                }

                for axis in chart_xml.findall(".//c:valAx", namespaces):
                    scaling = axis.find("c:scaling", namespaces)
                    if scaling is None:
                        scaling = parse_xml(f'<c:scaling xmlns:c="{namespaces["c"]}"/>')
                        axis.insert(0, scaling)

                    # Set min
                    min_elem = scaling.find("c:min", namespaces)
                    if min_elem is None:
                        min_elem = parse_xml(
                            f'<c:min xmlns:c="{namespaces["c"]}" val="{y_min}"/>'
                        )
                        scaling.append(min_elem)
                    else:
                        min_elem.set("val", str(y_min))

                    # Set max
                    max_elem = scaling.find("c:max", namespaces)
                    if max_elem is None:
                        max_elem = parse_xml(
                            f'<c:max xmlns:c="{namespaces["c"]}" val="{y_max}"/>'
                        )
                        scaling.append(max_elem)
                    else:
                        max_elem.set("val", str(y_max))
            except Exception as e:
                # Silently fail if auto-scaling doesn't work
                pass


def update_bar_chart_data(chart, categories, values):
    """Update bar chart data while preserving formatting"""
    chart_data = CategoryChartData()
    chart_data.categories = categories

    if chart.series:
        series_name = chart.series[0].name
    else:
        series_name = "Distribution"

    chart_data.add_series(series_name, values)
    chart.replace_data(chart_data)

    # CRITICAL: Re-apply formatting after replace_data()
    try:
        from pptx.util import Pt

        # Re-apply axis label font sizes and number format
        if hasattr(chart, 'category_axis') and chart.category_axis:
            chart.category_axis.tick_labels.font.size = Pt(9)
        if hasattr(chart, 'value_axis') and chart.value_axis:
            value_labels = chart.value_axis.tick_labels
            value_labels.font.size = Pt(9)
            # Use custom format to show % without multiplying by 100
            value_labels.number_format = '0"%"'
            value_labels.number_format_is_linked = False

        # Re-apply data labels
        if len(chart.plots) > 0:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(8)
            # Use custom format to show % with one decimal
            data_labels.number_format = '0.0"%"'
            data_labels.number_format_is_linked = False
    except Exception as e:
        # If formatting fails, at least data is correct
        pass


def extract_department_data(df, department):
    """Extract all data for a specific department"""

    dept_data = df[df["department"] == department]

    if dept_data.empty:
        raise ValueError(f"No data found for department: {department}")

    data = {}

    # Extract trend metrics
    for metric in ["satisfaction", "engagement", "response_rate", "retention"]:
        metric_row = dept_data[dept_data["metric_type"] == metric]
        if not metric_row.empty:
            row = metric_row.iloc[0]
            data[metric] = {
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "values": [row["Q1"], row["Q2"], row["Q3"], row["Q4"]],
            }

    # Extract age distribution
    age_row = dept_data[dept_data["metric_type"] == "age_distribution"]
    if not age_row.empty:
        row = age_row.iloc[0]
        data["age_distribution"] = {
            "categories": ["18-25", "26-35", "36-45", "46-55", "56+"],
            "values": [
                row["age_18_25"],
                row["age_26_35"],
                row["age_36_45"],
                row["age_46_55"],
                row["age_56_plus"],
            ],
        }

    # Extract summary info
    info_row = dept_data[dept_data["metric_type"] == "info"]
    if not info_row.empty:
        row = info_row.iloc[0]
        data["info"] = {
            "department_size": int(row["department_size"]),
            "latest_satisfaction": row["latest_satisfaction"],
            "latest_engagement": row["latest_engagement"],
            "latest_response_rate": row["latest_response_rate"],
        }

    return data


def find_shape_by_name(slide, shape_name):
    """Find shape by name (intelligent shape mapping)"""
    for shape in slide.shapes:
        if hasattr(shape, "name") and shape.name == shape_name:
            return shape
    return None


def find_shape_by_content(slide, search_text):
    """Find text shape by partial content match (fallback method)"""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text and search_text in shape.text:
            return shape
    return None


def find_chart_by_title(slide, chart_title_keyword):
    """Find chart by title keyword (fallback method)"""
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            if chart.has_title:
                title = chart.chart_title.text_frame.text.lower()
                if chart_title_keyword.lower() in title:
                    return shape
    return None


def update_slide(slide, department, dept_data):
    """
    Update all elements in slide for a department

    Uses intelligent shape mapping:
    1. Try to find shapes by name (from SHAPE_NAMES loaded via mapping.py)
    2. Fall back to content matching (slower, but works with unnamed shapes)
    """

    chart_count = 0

    # Update title
    title_shape = find_shape_by_name(slide, SHAPE_NAMES.get("title", "TextBox_Title"))
    if not title_shape:  # Fallback to content matching
        title_shape = find_shape_by_content(slide, "Quarterly Survey Report")
    if title_shape:
        update_text_preserve_format(
            title_shape, f"{department} - Quarterly Survey Report"
        )

    # Update subtitle
    subtitle_shape = find_shape_by_name(
        slide, SHAPE_NAMES.get("subtitle", "TextBox_Subtitle")
    )
    if not subtitle_shape:  # Fallback
        subtitle_shape = find_shape_by_content(slide, "Department Size")
    if subtitle_shape and "info" in dept_data:
        info = dept_data["info"]
        new_text = (
            f"Q4 2024 | Department Size: {info['department_size']} | "
            f"Response Rate: {info['latest_response_rate'] * 100:.1f}%"
        )
        update_text_preserve_format(subtitle_shape, new_text)

    # Update conclusion
    conclusion_shape = find_shape_by_name(
        slide, SHAPE_NAMES.get("conclusion", "TextBox_Conclusion")
    )
    if not conclusion_shape:  # Fallback
        conclusion_shape = find_shape_by_content(slide, "Key Findings")
    if conclusion_shape and "satisfaction" in dept_data:
        sat_values = dept_data["satisfaction"]["values"]
        sat_change = ((sat_values[-1] - sat_values[-2]) / sat_values[-2]) * 100
        direction = "increased" if sat_change > 0 else "decreased"

        new_text = (
            f"Key Findings: Overall satisfaction {direction} by {abs(sat_change):.1f}% QoQ. "
            f"Current satisfaction: {sat_values[-1] * 100:.1f}%, "
            f"engagement: {dept_data.get('engagement', {}).get('values', [0])[-1] * 100:.1f}%."
        )
        update_text_preserve_format(conclusion_shape, new_text)

    # Update charts using intelligent mapping
    for logical_name, title_keyword, data_key, num_format in CHART_MAPPINGS:
        # Try to find by name first
        chart_shape = find_shape_by_name(slide, SHAPE_NAMES.get(logical_name, ""))
        if not chart_shape:  # Fallback to title matching
            chart_shape = find_chart_by_title(slide, title_keyword)

        if chart_shape and chart_shape.has_chart and data_key in dept_data:
            data = dept_data[data_key]
            update_chart_data(
                chart_shape.chart,
                data["categories"],
                data["values"],
                number_format=num_format,
                auto_scale=True,  # Enable auto-scaling
            )
            chart_count += 1

    # Update Age Distribution chart
    age_chart_shape = find_shape_by_name(
        slide, SHAPE_NAMES.get("chart_age", "Chart_Age")
    )
    if not age_chart_shape:  # Fallback
        age_chart_shape = find_chart_by_title(slide, "Age")
    if (
        age_chart_shape
        and age_chart_shape.has_chart
        and "age_distribution" in dept_data
    ):
        data = dept_data["age_distribution"]
        update_bar_chart_data(age_chart_shape.chart, data["categories"], data["values"])
        chart_count += 1

    return chart_count


def update_detail_slide(slide, department, dept_data):
    """
    Update the detailed metrics slide (page 2) for a department.

    Fills:
    - Table_Metrics: Q3/Q4 + QoQ change for 4 key metrics
    - Table_Age: age distribution percentages
    - TextBox_Notes: short auto-generated notes
    """
    # Update title
    title_shape = find_shape_by_name(
        slide, SHAPE_NAMES.get("detail_title", "TextBox_DetailTitle")
    )
    if not title_shape:
        title_shape = find_shape_by_content(slide, "Detailed Metrics")
    if title_shape:
        update_text_preserve_format(title_shape, f"{department} - Detailed Metrics")

    # Subtitle can stay static or be updated if needed

    # Update key metrics table
    metrics_shape = find_shape_by_name(
        slide, SHAPE_NAMES.get("table_metrics", "Table_Metrics")
    )
    if metrics_shape and getattr(metrics_shape, "has_table", False):
        table = metrics_shape.table

        metrics_order = [
            ("satisfaction", "Satisfaction"),
            ("engagement", "Engagement"),
            ("response_rate", "Response Rate"),
            ("retention", "Retention"),
        ]

        for row_idx, (metric_key, display_name) in enumerate(metrics_order, start=1):
            if metric_key not in dept_data:
                continue

            values = dept_data[metric_key]["values"]
            # Use Q3/Q4 when available, otherwise last 2 points
            if len(values) >= 4:
                q3, q4 = values[2], values[3]
            elif len(values) >= 2:
                q3, q4 = values[-2], values[-1]
            else:
                continue

            if q3:
                delta_pp = (q4 - q3) * 100.0
            else:
                delta_pp = 0.0

            row = table.rows[row_idx]
            # Use format-preserving update for all cells
            update_table_cell_preserve_format(row.cells[0], display_name)
            update_table_cell_preserve_format(row.cells[1], f"{q3 * 100:.1f}%")
            update_table_cell_preserve_format(row.cells[2], f"{q4 * 100:.1f}%")
            update_table_cell_preserve_format(row.cells[3], f"{delta_pp:+.1f} pp")

    # Update age distribution table
    age_shape = find_shape_by_name(
        slide, SHAPE_NAMES.get("table_age", "Table_Age")
    )
    if (
        age_shape
        and getattr(age_shape, "has_table", False)
        and "age_distribution" in dept_data
    ):
        table = age_shape.table
        cats = dept_data["age_distribution"]["categories"]
        vals = dept_data["age_distribution"]["values"]

        for row_idx, (cat, val) in enumerate(zip(cats, vals), start=1):
            if row_idx >= len(table.rows):
                break
            row = table.rows[row_idx]
            # Use format-preserving update for all cells
            update_table_cell_preserve_format(row.cells[0], str(cat))
            update_table_cell_preserve_format(row.cells[1], f"{val:.1f}%")

    # Update Key Insights (consistent with page 1 style)
    notes_shape = find_shape_by_name(slide, SHAPE_NAMES.get("notes", "TextBox_Notes"))
    if not notes_shape:
        notes_shape = find_shape_by_content(slide, "Key Insights")
    if notes_shape and "satisfaction" in dept_data:
        sat_vals = dept_data["satisfaction"]["values"]
        if len(sat_vals) >= 2:
            q3, q4 = sat_vals[-2], sat_vals[-1]
            delta = (q4 - q3) * 100.0
            if delta > 0:
                direction = "improved"
            elif delta < 0:
                direction = "declined"
            else:
                direction = "stayed flat"

            # Build concise single-line insight (McKinsey style)
            insights = f"Key Insights: Satisfaction {direction} by {abs(delta):.1f} pp from Q3 to Q4."

            if "age_distribution" in dept_data:
                cats = dept_data["age_distribution"]["categories"]
                vals = dept_data["age_distribution"]["values"]
                if vals:
                    max_idx = max(range(len(vals)), key=lambda i: vals[i])
                    insights += f" Largest respondent group: {cats[max_idx]} ({vals[max_idx]:.1f}%)."

            update_text_preserve_format(notes_shape, insights)


def update_dashboard_slide(slide, department, dept_data):
    """
    Update the summary dashboard slide (page 1) for a department.

    Fills 6 business-style conclusion cards:
    - Card_OverallSatisfaction
    - Card_Engagement
    - Card_ResponseRetention
    - Card_AgeStructure
    - Card_SatVsEngagement
    - Card_OverallSummary
    """
    # Overall satisfaction
    card_sat = find_shape_by_name(slide, "Card_OverallSatisfaction")
    if card_sat and "satisfaction" in dept_data:
        vals = dept_data["satisfaction"]["values"]
        if len(vals) >= 2:
            q3, q4 = vals[-2], vals[-1]
            delta_pp = (q4 - q3) * 100.0
            text = (
                f"Overall Satisfaction\n"
                f"Q4: {q4 * 100:.1f}% ({delta_pp:+.1f} pp vs Q3)"
            )
        else:
            q4 = vals[-1]
            text = f"Overall Satisfaction\nQ4: {q4 * 100:.1f}%"
        update_text_preserve_format(card_sat, text)

    # Engagement
    card_eng = find_shape_by_name(slide, "Card_Engagement")
    if card_eng and "engagement" in dept_data:
        vals = dept_data["engagement"]["values"]
        if len(vals) >= 2:
            q3, q4 = vals[-2], vals[-1]
            delta_pp = (q4 - q3) * 100.0
            text = (
                f"Employee Engagement\n"
                f"Q4: {q4 * 100:.1f}% ({delta_pp:+.1f} pp vs Q3)"
            )
        else:
            q4 = vals[-1]
            text = f"Employee Engagement\nQ4: {q4 * 100:.1f}%"
        update_text_preserve_format(card_eng, text)

    # Response rate + retention
    card_resp_ret = find_shape_by_name(slide, "Card_ResponseRetention")
    if card_resp_ret:
        parts = []
        if "response_rate" in dept_data:
            vals = dept_data["response_rate"]["values"]
            q4 = vals[-1]
            parts.append(f"Response: {q4 * 100:.1f}%")
        if "retention" in dept_data:
            vals = dept_data["retention"]["values"]
            q4 = vals[-1]
            parts.append(f"Retention: {q4 * 100:.1f}%")
        if parts:
            text = "Response & Retention\n" + " | ".join(parts)
            update_text_preserve_format(card_resp_ret, text)

    # Age structure
    card_age = find_shape_by_name(slide, "Card_AgeStructure")
    if card_age and "age_distribution" in dept_data:
        cats = dept_data["age_distribution"]["categories"]
        vals = dept_data["age_distribution"]["values"]
        if vals:
            max_idx = max(range(len(vals)), key=lambda i: vals[i])
            main_group = cats[max_idx]
            share = vals[max_idx]
            text = (
                "Age Structure\n"
                f"Largest group: {main_group} ({share:.1f}%)"
            )
            update_text_preserve_format(card_age, text)

    # Satisfaction vs engagement gap
    card_gap = find_shape_by_name(slide, "Card_SatVsEngagement")
    if (
        card_gap
        and "satisfaction" in dept_data
        and "engagement" in dept_data
    ):
        sat_q4 = dept_data["satisfaction"]["values"][-1]
        eng_q4 = dept_data["engagement"]["values"][-1]
        gap_pp = (sat_q4 - eng_q4) * 100.0
        text = (
            "Satisfaction vs Engagement\n"
            f"Gap (Q4): {gap_pp:+.1f} pp"
        )
        update_text_preserve_format(card_gap, text)

    # Overall summary
    card_summary = find_shape_by_name(slide, "Card_OverallSummary")
    if card_summary and "satisfaction" in dept_data:
        sat_vals = dept_data["satisfaction"]["values"]
        if len(sat_vals) >= 2:
            q3, q4 = sat_vals[-2], sat_vals[-1]
            delta = (q4 - q3) * 100.0
            if delta > 0:
                direction = "improved"
            elif delta < 0:
                direction = "declined"
            else:
                direction = "stayed flat"
            summary = (
                f"Overall Summary\n"
                f"Satisfaction {direction} by {abs(delta):.1f} pp QoQ; "
            )
        else:
            summary = "Overall Summary\nSatisfaction stable QoQ; "

        if "age_distribution" in dept_data:
            cats = dept_data["age_distribution"]["categories"]
            vals = dept_data["age_distribution"]["values"]
            if vals:
                max_idx = max(range(len(vals)), key=lambda i: vals[i])
                summary += f"largest age group: {cats[max_idx]}."

        update_text_preserve_format(card_summary, summary)


def batch_generate(template_path, excel_path, output_dir):
    """Generate PPT reports for all departments"""

    print("\n" + "=" * 80)
    print("📊 Batch PPT Generation - Starting")
    print("=" * 80)

    # Load data
    print("\n📄 Loading data...")
    if not os.path.exists(excel_path):
        print(f"  ❌ Data file not found: {excel_path}")
        print(
            "  💡 Run 'python generate_sample_data.py' and "
            "'python prepare_data.py' first"
        )
        return

    df = pd.read_excel(excel_path)
    departments = df["department"].unique().tolist()
    print(f"  ✓ Found {len(departments)} departments")

    # Check template
    if not os.path.exists(template_path):
        print(f"  ❌ Template not found: {template_path}")
        print(f"  💡 Run 'python create_template.py' first")
        return

    print(f"  ✓ Template loaded")

    # Create output directory
    os.makedirs(output_dir, exist_ok=True)

    # Generate reports
    print("\n" + "=" * 80)
    print("📌 Generating Reports")
    print("=" * 80)

    results = {"success": [], "failed": []}

    for idx, department in enumerate(departments, 1):
        try:
            print(f"\n[{idx}/{len(departments)}] {department}...")

            # Load fresh template
            prs = Presentation(template_path)

            # Extract department data
            dept_data = extract_department_data(df, department)

            # Update dashboard slide (page 1), if present
            if len(prs.slides) > 0:
                update_dashboard_slide(prs.slides[0], department, dept_data)
                print("  ✓ Updated summary dashboard")

            # Update overview slide (page 2)
            chart_count = 0
            if len(prs.slides) > 1:
                chart_count = update_slide(prs.slides[1], department, dept_data)
            print(f"  ✓ Updated {chart_count} charts on overview slide")

            # Update detail slide (page 3), if present
            if len(prs.slides) > 2:
                update_detail_slide(prs.slides[2], department, dept_data)
                print("  ✓ Updated detail tables")

            # Save
            output_path = os.path.join(output_dir, f"{department}_Report.pptx")
            prs.save(output_path)
            print(f"  ✓ Saved: {output_path}")

            results["success"].append(department)

        except Exception as e:
            print(f"  ❌ Failed: {e}")
            results["failed"].append((department, str(e)))

    # Summary
    print("\n" + "=" * 80)
    print("📊 Generation Complete")
    print("=" * 80)
    print(f"\n✅ Success: {len(results['success'])} reports")
    print(f"❌ Failed: {len(results['failed'])} reports")

    if results["success"]:
        print("\n📁 Generated reports:")
        for dept in results["success"]:
            print(f"  - {dept}_Report.pptx")

    if results["failed"]:
        print("\n⚠️  Failed reports:")
        for dept, error in results["failed"]:
            print(f"  - {dept}: {error}")

    print("\n💡 All reports saved to: " + output_dir)
    print("=" * 80 + "\n")

    return results


def main():
    """Entry point"""
    template_path = "data/raw/template.pptx"
    excel_path = "data/processed/department_metrics.xlsx"
    output_dir = "data/output/reports"

    try:
        results = batch_generate(template_path, excel_path, output_dir)

        if results and len(results["success"]) > 0:
            print("🎉 Batch generation successful!")
            print(f"   Generated {len(results['success'])} reports in seconds")
            print(f"   vs. hours of manual work!")

    except Exception as e:
        print(f"\n❌ Batch generation failed: {e}")
        import traceback

        traceback.print_exc()


if __name__ == "__main__":
    main()
