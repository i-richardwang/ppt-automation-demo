#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT Template Analyzer

Analyzes PPT template structure to help understand shape names and indices.
"""

from pptx import Presentation
import sys


def analyze_template(ppt_path):
    """
    Analyze PPT template structure

    Args:
        ppt_path: Path to PPT file
    """
    print(f"\n{'='*80}")
    print(f"PPT Template Analysis: {ppt_path}")
    print("="*80)

    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"❌ Cannot open file: {e}")
        return

    print(f"\nTotal slides: {len(prs.slides)}\n")

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nSlide {slide_idx}:")
        print("-"*80)

        shape_count = len(slide.shapes)
        chart_count = sum(1 for s in slide.shapes if s.has_chart)
        text_count = sum(1 for s in slide.shapes if s.has_text_frame)

        print(f"  Total shapes: {shape_count} (Charts: {chart_count}, Text: {text_count})")
        print()

        for shape_idx, shape in enumerate(slide.shapes):
            # Determine shape type
            if shape.has_chart:
                shape_type = "Chart"
                chart = shape.chart
                series_count = len(chart.series)
                series_names = [s.name for s in chart.series]

                # Get chart title
                chart_title = ""
                if chart.has_title:
                    chart_title = chart.chart_title.text_frame.text

                extra_info = f", Series: {series_count}, Names: {series_names}"
                if chart_title:
                    extra_info += f", Title: \"{chart_title}\""

            elif shape.has_text_frame:
                shape_type = "Text"
                text = shape.text.strip()
                text_preview = text[:40] + "..." if len(text) > 40 else text
                extra_info = f", Content: \"{text_preview}\""

            elif shape.has_table:
                shape_type = "Table"
                table = shape.table
                extra_info = f", Size: {len(table.rows)}×{len(table.columns)}"

            else:
                shape_type = "Other"
                extra_info = ""

            print(f"  [{shape_idx:2d}] {shape.name:30s} ({shape_type}){extra_info}")

    print("\n" + "="*80)
    print("Analysis Complete!")
    print("="*80)
    print("\n💡 Tips:")
    print("  - Use PowerPoint's 'Selection Pane' to rename shapes")
    print("  - Recommended naming: Chart_1, Chart_2, TextBox_Title, etc.")
    print("  - Charts must have pre-set series names for proper legend display")
    print("="*80)


def main():
    """Entry point"""
    if len(sys.argv) < 2:
        print("Usage: python analyze_template.py <ppt_path>")
        print("\nExample:")
        print("  python tools/analyze_template.py data/raw/template.pptx")
        sys.exit(1)

    ppt_path = sys.argv[1]
    analyze_template(ppt_path)


if __name__ == "__main__":
    main()
